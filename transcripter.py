import os
import json
import logging
import subprocess
 
from agno.agent import Agent
from agno.models.groq import Groq as Agno_Groq

from dotenv import load_dotenv
from groq import Groq
from pydub import AudioSegment
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from PIL import Image as PILImage
import ffmpeg
 
# ---------------------------------------------------------------------------
# Logging setup
# ---------------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s — %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger("video_study_gen")
 
load_dotenv()
 
 
# ---------------------------------------------------------------------------
# Environment / pre-flight validation
# ---------------------------------------------------------------------------
 
def validate_env() :
    """
    Check all required environment variables and external tools before running.
    Returns True if everything is OK, False otherwise.
    Logs a specific error for each missing requirement.
    """
    ok = True
 
    # --- API keys ---
    if not os.getenv("GROQ_API_KEY"):
        logger.error(
            "GROQ_API_KEY is not set. "
            "Add it to your .env file or export it as an environment variable."
        )
        ok = False
 
    # --- Required files ---
    for fname in ("quiz.md", "slide.md"):
        if not os.path.isfile(fname):
            logger.error("Required prompt file not found: %s", fname)
            ok = False
 
    # --- External CLI tools ---
    for tool in ("yt-dlp", "ffmpeg"):
        result = subprocess.run(
            ["which", tool], capture_output=True, text=True
        )
        if result.returncode != 0:
            logger.error(
                "Required CLI tool '%s' not found. "
                "Install it and make sure it is on your PATH.",
                tool,
            )
            ok = False
 
    # --- Frames directory ---
    frames_dir = Video_decompose.frames_dir
    if not os.path.isdir(frames_dir):
        try:
            os.makedirs(frames_dir, exist_ok=True)
            logger.info("Created frames directory: %s", frames_dir)
        except OSError as exc:
            logger.error("Could not create frames directory '%s': %s", frames_dir, exc)
            ok = False
 
    return ok
 
 
# ---------------------------------------------------------------------------
# Video decomposition
# ---------------------------------------------------------------------------
 
class Video_decompose:
    output = "video.mp4"
    frames_dir = "frames"
 
    def download_from_url(self, url: str) :
        logger.info("Downloading video from URL: %s", url)
        result = subprocess.run(
            ["yt-dlp", url, "-o", self.output],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            logger.error(
                "yt-dlp failed (exit code %d).\nstdout: %s\nstderr: %s",
                result.returncode,
                result.stdout,
                result.stderr,
            )
            return
        logger.info("Download complete → %s", self.output)
 
    def split_audio(self) -> list[str]:
        webm_path = self.output
        if os.path.exists(webm_path + ".webm") : 
            webm_path += ".webm"
            self.output += ".webm"
        if not os.path.exists(webm_path):
            logger.error("Audio file not found after download: %s", webm_path)
            return []
        logger.info("Splitting audio into 60-second chunks…")
        chunk_duration_ms = 60_000
        audio = AudioSegment.from_file(webm_path)
        chunks: list[str] = []
 
        for i in range(0, len(audio), chunk_duration_ms):
            chunk = audio[i : i + chunk_duration_ms]
            chunk_path = f"/tmp/chunk_{i}.mp3"
            chunk.export(chunk_path, format="mp3", bitrate="64k")
            chunks.append(chunk_path)
            logger.debug("Exported chunk: %s", chunk_path)
 
        logger.info("Audio split into %d chunk(s).", len(chunks))
        return chunks
 
    def transcript_chunks(self, chunks: list[str]) -> str:
        if not chunks:
            logger.warning("No audio chunks to transcribe.")
            return ""
 
        client = Groq()
        full_transcript = ""
 
        for chunk_path in chunks:
            logger.info("Transcribing: %s", chunk_path)
            try:
                with open(chunk_path, "rb") as f:
                    result = client.audio.transcriptions.create(
                        model="whisper-large-v3",
                        file=f,
                    )
                full_transcript += result.text + " "
                os.remove(chunk_path)
                logger.debug("Removed temp chunk: %s", chunk_path)
 
            except Groq.AuthenticationError:
                logger.error(
                    "Invalid Groq API key. Check GROQ_API_KEY in your .env file."
                )
            except Groq.RateLimitError:
                logger.error(
                    "Groq rate limit or token quota exceeded. "
                    "Wait a moment or upgrade your plan."
                )
            except Groq.APIConnectionError as exc:
                logger.error("Could not reach Groq API: %s", exc)
            except Groq.GroqError as exc:
                logger.error("Groq API error on chunk '%s': %s", chunk_path, exc)
 
        transcript = full_transcript.strip()
        logger.info(
            "Transcription complete. Total characters: %d", len(transcript)
        )
        return transcript
 
    def extract_frames(self, n_frames: int) :
        webm_path = self.output
        if not os.path.exists(webm_path):
            logger.error("Cannot extract frames — file not found: %s", webm_path)
            return
 
        audio = AudioSegment.from_file(webm_path)
        fps = 1000 * n_frames / len(audio)
        logger.info(
            "Extracting %d frame(s) at %.4f fps from %s…", n_frames, fps, webm_path
        )
 
        try:
            (
                ffmpeg
                .input(webm_path)
                .filter("fps", fps)
                .output(f"{self.frames_dir}/frame_%04d.jpg")
                .run(quiet=True)
            )
            logger.info("Frames saved to '%s/'.", self.frames_dir)
        except ffmpeg.Error as exc:
            logger.error("ffmpeg frame extraction failed: %s", exc.stderr.decode())
 
    def transcribe_large_audio(self, url: str) -> str:
        self.download_from_url(url)
        chunks = self.split_audio()
        return self.transcript_chunks(chunks)
 
 
# ---------------------------------------------------------------------------
# Agent tools
# ---------------------------------------------------------------------------
 
class Agent_tools:
    agent = Agent(model=Agno_Groq(id="openai/gpt-oss-120b"))
    quiz_prompt = "quiz.md"
    slide_prompt = "slide.md"
    frames_folder = Video_decompose.frames_dir
 
    def file_to_prompt(self, filename: str) -> str:
        try:
            content = open(filename, "r", encoding="utf-8").read()
            logger.debug("Loaded prompt from '%s' (%d chars).", filename, len(content))
            return content
        except OSError as exc:
            logger.error("Could not read prompt file '%s': %s", filename, exc)
            return ""
 
    def build_pptx(self, slides: list[dict]) :
        filename = "presentation.pptx"
        logger.info("Building PPTX with %d slide(s)…", len(slides))
 
        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        half_width = slide_width // 2
        blank_layout = prs.slide_layouts[6]
 
        image_files = sorted(os.listdir(self.frames_folder))
        if len(image_files) < len(slides):
            logger.warning(
                "Only %d frame image(s) available for %d slide(s). "
                "Some slides may be missing images.",
                len(image_files),
                len(slides),
            )
 
        for i, slide_data in enumerate(slides):
            slide = prs.slides.add_slide(blank_layout)
 
            # Black left-half background
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=Emu(0),
                top=Emu(0),
                width=half_width,
                height=slide_height,
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            bg.line.fill.background()
 
            # Title
            title_top = slide_height * 15 // 100
            title_height = slide_height * 15 // 100
            title_box = slide.shapes.add_textbox(
                left=Emu(360_000),
                top=title_top,
                width=half_width - Emu(720_000),
                height=title_height,
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", "")
            run = p.runs[0]
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.bold = True
            run.font.size = Pt(28)
 
            # Bullets
            bullets_top = title_top + title_height + Emu(200_000)
            bullets_box = slide.shapes.add_textbox(
                left=Emu(360_000),
                top=bullets_top,
                width=half_width - Emu(720_000),
                height=slide_height - bullets_top - Emu(360_000),
            )
            btf = bullets_box.text_frame
            btf.word_wrap = True
            btf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
 
            for j, bullet in enumerate(slide_data.get("bullets", [])):
                p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = Pt(6)
                run = p.runs[0]
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(16)
 
            # Right-half image
            if i < len(image_files):
                image_path = os.path.join(self.frames_folder, image_files[i])
                try:
                    with PILImage.open(image_path) as img:
                        img_w, img_h = img.size
 
                    scale = min(half_width / img_w, slide_height / img_h)
                    pic_w = int(img_w * scale)
                    pic_h = int(img_h * scale)
                    pic_left = half_width + (half_width - pic_w) // 2
                    pic_top = (slide_height - pic_h) // 2
 
                    slide.shapes.add_picture(
                        image_path,
                        left=pic_left,
                        top=pic_top,
                        width=pic_w,
                        height=pic_h,
                    )
                    logger.debug("Added image '%s' to slide %d.", image_path, i + 1)
                except (OSError, Exception) as exc:
                    logger.warning(
                        "Could not add image '%s' to slide %d: %s",
                        image_path, i + 1, exc,
                    )
            else:
                logger.warning("No image available for slide %d.", i + 1)
 
            # Speaker notes
            slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
 
        prs.save(filename)
        logger.info("Presentation saved → %s", filename)
 
    def gen_questions(self, url: str) :
        logger.info("Generating quiz questions for: %s", url)
        prompt = self.file_to_prompt(self.quiz_prompt)
        if not prompt: return
        Vc = Video_decompose()
        transcript = Vc.transcribe_large_audio(url)
        if not transcript:
            logger.error("Empty transcript — cannot generate questions.")
            return
        os.remove(Vc.output)
        logger.info("Deleted temp file: %s", Vc.output)
        for frame in os.listdir(Vc.frames_dir):
            os.remove(os.path.join(Vc.frames_dir, frame))
        self.agent.print_response(prompt + "\n" + transcript)
 
    def gen_slides(self, url: str) :
        logger.info("Generating slides for: %s", url)
        prompt = self.file_to_prompt(self.slide_prompt)
        Vc = Video_decompose()
        if not prompt: return
        transcript = Video_decompose().transcribe_large_audio(url)
        if not transcript:
            logger.error("Empty transcript — cannot generate slides.")
            return
 
        raw = self.agent.run(prompt + "\n" + transcript).content
        try:
            slides = json.loads(raw.strip().strip("```json").strip("```"))
        except json.JSONDecodeError as exc:
            logger.error("Failed to parse slides JSON from model response: %s", exc)
            logger.debug("Raw model response:\n%s", raw)
            return

        Video_decompose().extract_frames(len(slides))
        self.build_pptx(slides)
        os.remove(Vc.output)
        logger.info("Deleted temp file: %s", Vc.output)
        for frame in os.listdir(Vc.frames_dir):
            os.remove(os.path.join(Vc.frames_dir, frame))

 
 
# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------
 
if __name__ == "__main__":
    print("---------WELCOME!----------")
 
    if not validate_env():
        logger.critical(
            "Pre-flight checks failed. Fix the errors above and try again."
        )
        raise SystemExit(1)
 
    url = input("URL: ").strip()
    if not url:
        logger.error("No URL provided. Exiting.")
        raise SystemExit(1)
 
    print("Digit:")
    print("0 - Generate video quiz")
    print("1 - Generate slide")
    mode = input().strip()
 
    AT = Agent_tools()
    if mode == "0":
        AT.gen_questions(url)
    elif mode == "1":
        AT.gen_slides(url)
    else:
        logger.error("Invalid mode '%s'. Choose 0 or 1.", mode)
        raise SystemExit(1)
 
